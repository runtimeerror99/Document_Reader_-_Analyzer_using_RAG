import streamlit as st
import os
import pptx
import pyrebase
import time
from menu import menu
from llama_index.core import SimpleDirectoryReader, VectorStoreIndex, SummaryIndex

os.environ["OPENAI_API_KEY"] = st.secrets["openai"]

firebaseConfig = {
    'apiKey': st.secrets["apiKey"],
    'authDomain': st.secrets["authDomain"],
    'projectId': st.secrets["projectId"],
    'storageBucket': st.secrets["storageBucket"],
    'messagingSenderId': st.secrets["messagingSenderId"],
    'appId': st.secrets["appId"],
    'measurementId': st.secrets["measurementId"],
    'databaseURL': st.secrets["databaseURL"]
}

firebase = pyrebase.initialize_app(firebaseConfig)

st.set_page_config(page_title="DORA", page_icon="🦙")
st.markdown(f"""<style>
        .st-emotion-cache-79elbk{{
            display: none;}}
            </style>""", unsafe_allow_html=True)
menu()

if "role" not in st.session_state or st.session_state.role is None:
    st.switch_page('./pages/authenticate.py')

if "projects" not in st.session_state:
    st.session_state.projects = []

if "curr" not in st.session_state:
    st.session_state.curr = None


def create_project(project_name):
    if project_name not in st.session_state.projects:
        st.session_state.projects.append(project_name)
        st.success(f"Project '{project_name}' created successfully.")


def upload_and_store_files(file_uploads, target_folder):
    file_paths = []
    if not os.path.exists(target_folder):
        os.makedirs(target_folder)

    for file_upload in file_uploads:
        file_path = os.path.join(target_folder, file_upload.name)
        
        if file_path.endswith(".pptx"):
            ppt = pptx.Presentation(file_upload)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            with open(f'{target_folder}/{str(file_upload.name).split(".")[0]}.txt', "w") as f:
                f.write("\n".join(text))
        else:
            with open(file_path, "wb") as f:
                f.write(file_upload.getbuffer())
        
        file_paths.append(file_path)

    return file_paths
def show_files(project_name):
    if os.path.exists(f"{st.session_state.role}/{project_name}"):
        files = os.listdir(f"{st.session_state.role}/{project_name}")
        if files:
            for file in files:
                st.write('--'+file)
        else:
            st.write("The project is empty.")



def create_index(project_name):
    docs = SimpleDirectoryReader(f"{st.session_state.role}/{project_name}").load_data()
    index = VectorStoreIndex.from_documents(docs)
    summary = SummaryIndex.from_documents(docs)

    index.storage_context.persist(f'{st.session_state.role}/index/{project_name}')
    summary.storage_context.persist(f'{st.session_state.role}/summary/{project_name}')


st.header("Select or Create Project")
project_name = st.selectbox("Select or Create Project:", options=st.session_state.projects + ["Create New Project"])

if project_name == "Create New Project":
    new_project_name = st.text_input("Enter New Project Name:")

uploaded_files = st.file_uploader("Upload multiple files", accept_multiple_files=True)

process_button = st.button("Process Project 🚀")

if process_button:
    if project_name == "Create New Project":
        if not new_project_name:
            st.error("Please enter a project name.")
        else:
            project_name = new_project_name
            create_project(project_name)

    if project_name:
        target_folder = f"{st.session_state.role}/{project_name}/"
        
        with st.spinner("Uploading files... 📂"):
            file_paths = upload_and_store_files(uploaded_files, target_folder)
            if file_paths:
                st.success("Files uploaded and stored successfully.")
            else:
                st.warning("No files uploaded.")

        with st.spinner("Creating index... ⏳"):
            create_index(project_name)
            st.success("Index created! Navigate to the Query page to start querying.")
            st.toast("Index Creation Successful!", icon="🎉")

st.write("Files in this project:")
if project_name != "Create New Project":
    show_files(project_name)
